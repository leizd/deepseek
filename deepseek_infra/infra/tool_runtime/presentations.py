"""Generate downloadable .pptx decks from a structured outline.

Exposed to the model as the ``create_pptx`` function-calling tool: the model
passes a title and a list of slides, we render a real PowerPoint file into
``.generated/{id}.pptx`` and hand back a ``downloadUrl`` the model surfaces to
the user as a Markdown link. No frontend changes needed — the chat renderer
already turns Markdown links into clickable anchors.
"""

from __future__ import annotations

import re
from typing import Any

from deepseek_infra.core.errors import AppError, ErrorCode
from deepseek_infra.infra.tool_runtime.generated_files import (
    cleanup_generated_files,
    resolve_generated_file,
    save_generated_file_to_downloads,
    store_generated_file,
)
from deepseek_infra.infra.tool_runtime.slides_skill import SLIDES_SKILL_NAME

# resolve/save/cleanup 现在由 generated_files 统一提供，这里 re-export 以保持历史导入路径。
__all__ = [
    "cleanup_generated_files",
    "create_presentation",
    "create_presentation_from_text",
    "infer_presentation_title",
    "resolve_generated_file",
    "save_generated_file_to_downloads",
    "slides_from_outline_text",
]

MAX_SLIDES = 60
MAX_BULLETS_PER_SLIDE = 24
_OUTLINE_HEADING_RE = re.compile(
    r"^\s*(?:#{1,4}\s*)?(?:\*\*)?\s*"
    r"(?:(?:第\s*)?\d{1,2}\s*(?:页|张)?\s*[、.．):：-]|(?:幻灯片|页面|页|slide)\s*\d{1,2}\s*[、.．):：-]?)"
    r"\s*(.+?)\s*(?:\*\*)?\s*$",
    flags=re.IGNORECASE,
)
_MARKDOWN_SLIDE_HEADING_RE = re.compile(r"^\s*#{2,4}\s+(.+?)\s*$")
_BULLET_RE = re.compile(
    r"^\s*(?:[-*•·]\s+|[（(]?[a-zA-Z][）).、]\s*|[（(]?\d{1,2}[）)]\s*|\d{1,2}[、.．)]\s*)(.+?)\s*$"
)
_REFUSAL_RE = re.compile(r"(?:无法|不能|没有.*能力|只能|可以帮.*大纲|不能直接|无法直接)")
_OUTLINE_META_TITLE_RE = re.compile(r"(?:ppt|powerpoint|presentation|幻灯片|演示文稿).{0,12}大纲|^大纲$", re.IGNORECASE)


def _normalize_bullets(item: dict[str, Any]) -> list[str]:
    bullets = item.get("bullets")
    if isinstance(bullets, list):
        return [str(b).strip() for b in bullets if str(b).strip()]
    content = str(item.get("content") or "").strip()
    if content:
        return [line.strip() for line in content.splitlines() if line.strip()]
    return []


def infer_presentation_title(prompt: str, outline: str = "") -> str:
    text = re.sub(r"\s+", " ", str(prompt or "")).strip()
    patterns = (
        r"(?:关于|介绍|讲解|主题为|题为)\s*([A-Za-z0-9\u4e00-\u9fff _-]{1,48}?)\s*(?:的)?\s*(?:PPT|ppt|幻灯片|演示文稿|presentation)",
        r"(?:PPT|ppt|幻灯片|演示文稿|presentation).*?(?:关于|介绍|讲解)\s*([A-Za-z0-9\u4e00-\u9fff _-]{1,48})",
        r"([A-Za-z0-9\u4e00-\u9fff _-]{1,48}?)\s*(?:PPT|ppt|幻灯片|演示文稿|presentation)",
    )
    for pattern in patterns:
        match = re.search(pattern, text, flags=re.IGNORECASE)
        if match:
            title = _clean_outline_title(match.group(1))
            title = re.sub(r"^(?:一个|一份|一套|做一个|做一份|制作一个|制作一份|帮我|请帮我)\s*", "", title)
            if title:
                normalized = _normalize_topic_name(title)
                if re.search(r"\b介绍\b|介绍", text) and "介绍" not in normalized:
                    return f"{normalized} 介绍"
                return normalized

    for line in str(outline or "").splitlines():
        cleaned = _clean_outline_title(line)
        if cleaned and len(cleaned) <= 60 and re.search(r"(?:PPT|幻灯片|演示文稿|介绍)$", cleaned, flags=re.IGNORECASE):
            return _normalize_topic_name(cleaned)
    return "演示文稿"


def create_presentation_from_text(prompt: str, outline: str) -> dict[str, Any]:
    title = infer_presentation_title(prompt, outline)
    slides = slides_from_outline_text(outline, topic=title)
    return create_presentation(title, slides, subtitle=f"由 DeepSeek Infra {SLIDES_SKILL_NAME} skill 本地生成")


def slides_from_outline_text(outline: str, *, topic: str = "") -> list[dict[str, Any]]:
    text = str(outline or "").replace("\r\n", "\n").replace("\r", "\n")
    text = _outline_relevant_tail(text)
    slides: list[dict[str, Any]] = []
    current: dict[str, Any] | None = None
    for raw_line in text.splitlines():
        line = raw_line.strip()
        if not line or line.startswith("```"):
            continue
        if current is not None and _looks_like_numbered_body_line(line):
            bullet = _BULLET_RE.match(line)
            if bullet:
                text = _clean_outline_bullet(bullet.group(1))
                if text and not _REFUSAL_RE.search(text):
                    current["bullets"].append(text)
                continue
        title = _outline_slide_title(line, has_slides=bool(slides))
        if title is not None:
            if not title or _REFUSAL_RE.search(title):
                continue
            current = {"title": title, "bullets": []}
            slides.append(current)
            continue
        if current is None:
            continue
        bullet = _BULLET_RE.match(line)
        if bullet:
            text = _clean_outline_bullet(bullet.group(1))
            if text and not _REFUSAL_RE.search(text):
                current["bullets"].append(text)
        elif raw_line.startswith((" ", "\t")):
            text = _clean_outline_bullet(line)
            if text and not _REFUSAL_RE.search(text):
                current["bullets"].append(text)
        else:
            text = _clean_outline_bullet(line)
            if _looks_like_body_line(text) and not _REFUSAL_RE.search(text):
                current["bullets"].append(text)

    slides = _drop_duplicate_cover_slide(slides, topic)
    result: list[dict[str, Any]] = []
    for slide in slides:
        title = _clean_outline_title(slide.get("title", ""))
        if not title:
            continue
        bullets = [str(item).strip() for item in slide.get("bullets", []) if str(item).strip()]
        if not bullets:
            bullets = _default_bullets_for_title(title, topic)
        result.append({"title": title, "bullets": bullets[:MAX_BULLETS_PER_SLIDE]})
        if len(result) >= MAX_SLIDES:
            break
    return result or _default_slides(topic)


def _outline_relevant_tail(text: str) -> str:
    markers = ("PPT 大纲", "ppt 大纲", "幻灯片大纲", "演示文稿大纲", "关于")
    best = -1
    lowered = text.lower()
    for marker in markers:
        index = lowered.rfind(marker.lower())
        if index > best:
            best = index
    return text[best:] if best >= 0 else text


def _outline_slide_title(line: str, *, has_slides: bool) -> str | None:
    heading = _OUTLINE_HEADING_RE.match(line)
    if heading:
        return _clean_outline_title(heading.group(1))

    markdown_heading = _MARKDOWN_SLIDE_HEADING_RE.match(line)
    if not markdown_heading:
        return None
    title = _clean_outline_title(markdown_heading.group(1))
    if not has_slides and _OUTLINE_META_TITLE_RE.search(title):
        return None
    return title


def _clean_outline_title(value: object) -> str:
    text = str(value or "").strip().strip("#*`-—:：.。 ")
    text = re.sub(r"^(?:标题|页面|幻灯片|Slide)\s*[:：-]\s*", "", text, flags=re.IGNORECASE)
    text = re.sub(r"^封面\s*[-:：]\s*", "", text)
    return text[:120].strip()


def _clean_outline_bullet(value: object) -> str:
    text = str(value or "").strip().strip("-*•· ")
    return text[:500].strip()


def _looks_like_body_line(text: str) -> bool:
    if not text:
        return False
    if _OUTLINE_META_TITLE_RE.search(text):
        return False
    if text.startswith(("```", "|")):
        return False
    return len(text) <= 240


def _looks_like_numbered_body_line(line: str) -> bool:
    return bool(re.match(r"^\s*(?:\d{1,2}[、)]|[（(]\d{1,2}[）)])", str(line or "")))


def _normalize_topic_name(value: str) -> str:
    text = re.sub(r"\s+", " ", str(value or "")).strip(" ：:-")
    if text.lower() == "git":
        return "Git 介绍"
    return text[:80] or "演示文稿"


def _drop_duplicate_cover_slide(slides: list[dict[str, Any]], topic: str) -> list[dict[str, Any]]:
    if len(slides) <= 1:
        return slides
    first_title = str(slides[0].get("title") or "")
    compact_topic = re.sub(r"\s+", "", topic.lower())
    compact_first = re.sub(r"\s+", "", first_title.lower())
    if "封面" in first_title or (compact_topic and compact_first in {compact_topic, f"{compact_topic}介绍"}):
        return slides[1:]
    return slides


def _default_bullets_for_title(title: str, topic: str) -> list[str]:
    topic_text = topic or "本主题"
    return [
        f"说明“{title}”在 {topic_text} 中的作用",
        "列出关键概念、常见场景和易错点",
        "结合简短示例帮助听众快速理解",
    ]


def _default_slides(topic: str) -> list[dict[str, Any]]:
    normalized_topic = topic or "主题"
    if re.search(r"\bgit\b", normalized_topic, flags=re.IGNORECASE):
        return [
            {"title": "什么是 Git", "bullets": ["分布式版本控制系统", "记录代码和文档的每次变更", "支持多人并行协作"]},
            {"title": "为什么需要版本控制", "bullets": ["追踪历史，方便回滚", "多人协作时降低冲突成本", "用分支隔离不同功能和实验"]},
            {"title": "核心概念", "bullets": ["仓库 Repository", "提交 Commit", "分支 Branch", "合并 Merge", "远程仓库 Remote"]},
            {"title": "基本工作流程", "bullets": ["clone 或 init 创建仓库", "add 暂存变更", "commit 保存快照", "push / pull 同步远程"]},
            {"title": "常用命令", "bullets": ["git status 查看状态", "git add / commit 提交变更", "git log 查看历史", "git branch / merge 管理分支"]},
            {"title": "分支与协作", "bullets": ["为功能创建独立分支", "通过 Pull Request / Merge Request 审查", "解决冲突后再合并到主线"]},
            {"title": "最佳实践", "bullets": ["小步提交并写清楚提交信息", "不要提交密钥和大文件", "合并前先同步主分支", "用 .gitignore 排除临时文件"]},
            {"title": "总结", "bullets": ["Git 让变更可追踪、协作可控", "掌握工作流比死记命令更重要", "从日常 add / commit / push 开始练习"]},
        ]
    return [
        {"title": f"{normalized_topic}概览", "bullets": ["背景与目标", "核心问题", "本次演示的主要结论"]},
        {"title": "核心概念", "bullets": ["关键术语", "基本原理", "适用范围"]},
        {"title": "方法与流程", "bullets": ["准备工作", "执行步骤", "检查与反馈"]},
        {"title": "案例或示例", "bullets": ["典型场景", "操作示例", "结果解读"]},
        {"title": "风险与建议", "bullets": ["常见误区", "注意事项", "改进建议"]},
        {"title": "总结", "bullets": ["回顾重点", "下一步行动", "问答"]},
    ]


# 一组现代配色主题；按标题哈希确定性选一套，让同一份内容稳定、不同主题之间有区分。
# 每套：cover_bg 封面深色底，accent 强调色（内容页标题/线/bullet 符号），band 封面副标题浅色。
_DECK_THEMES = (
    {"cover_bg": "1E3A8A", "accent": "2563EB", "band": "BFDBFE"},  # 蓝
    {"cover_bg": "312E81", "accent": "6366F1", "band": "C7D2FE"},  # 靛紫
    {"cover_bg": "0F766E", "accent": "0D9488", "band": "99F6E4"},  # 青绿
    {"cover_bg": "9A3412", "accent": "EA580C", "band": "FED7AA"},  # 暖橙
    {"cover_bg": "0F172A", "accent": "0EA5E9", "band": "BAE6FD"},  # 深空蓝
    {"cover_bg": "581C87", "accent": "A855F7", "band": "E9D5FF"},  # 紫
)
_TITLE_INK = "111827"   # 标题近黑
_BODY_INK = "334155"    # 正文深灰
_DETAIL_INK = "64748B"  # 次级说明文字
_MUTED_INK = "94A3B8"   # 页码 / 次要
_HAIRLINE = "E2E8F0"    # 细分隔线
_CN_FONT = "微软雅黑"


def _hex(value: str) -> Any:
    from pptx.dml.color import RGBColor

    return RGBColor.from_string(value)


def _deck_theme_for(title: str) -> dict[str, str]:
    import hashlib

    digest = hashlib.md5(str(title or "").encode("utf-8")).hexdigest()
    return _DECK_THEMES[int(digest, 16) % len(_DECK_THEMES)]


def _set_slide_bg(slide: Any, color: str) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _hex(color)


def _rule(slide: Any, x: float, y: float, w: float, h: float, color: str) -> None:
    """画一条无边框、无阴影的实心细条/方块，用于强调线、分隔线和小标记。"""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _hex(color)
    bar.line.fill.background()
    bar.shadow.inherit = False


def _add_cover_slide(prs: Any, theme: dict[str, str], title: str, subtitle: str) -> None:
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, theme["cover_bg"])

    # 左侧细强调竖条 + 标题上方 eyebrow 短线
    _rule(slide, 0.0, 0.0, 0.3, 7.5, theme["accent"])
    _rule(slide, 1.2, 2.35, 0.7, 0.06, theme["accent"])

    box = slide.shapes.add_textbox(Inches(1.2), Inches(2.55), Inches(10.9), Inches(2.6))
    tf = box.text_frame
    tf.word_wrap = True
    _style_paragraph(tf.paragraphs[0], title[:120], size=(44 if len(title) > 18 else 52), color="FFFFFF", bold=True)
    if subtitle:
        sp = tf.add_paragraph()
        _style_paragraph(sp, subtitle[:160], size=18, color=theme["band"])
        sp.space_before = Pt(18)


def _add_content_slide(prs: Any, theme: dict[str, str], title: str, bullets: list[str], page_no: int) -> None:
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, "FFFFFF")
    _add_header(slide, theme, title)

    # 开放式要点列表：每条 lead 加粗、detail 次级灰，行间细线分隔，无项目符号方块
    count = max(1, len(bullets))
    lead_size = 19 if count <= 4 else 17 if count <= 6 else 15
    detail_size = max(11, lead_size - 5)
    top, avail = 2.0, 4.9
    step = avail / count
    for index, bullet in enumerate(bullets):
        y = top + index * step
        lead, detail = _split_lead_detail(bullet)
        _rule(slide, 0.9, y + 0.07, 0.1, 0.1, theme["accent"])
        box = slide.shapes.add_textbox(Inches(1.2), Inches(y), Inches(11.2), Inches(step))
        tf = box.text_frame
        tf.word_wrap = True
        _style_paragraph(tf.paragraphs[0], lead[:240], size=lead_size, color=_TITLE_INK, bold=True)
        if detail:
            dp = tf.add_paragraph()
            _style_paragraph(dp, detail[:300], size=detail_size, color=_DETAIL_INK)
            dp.space_before = Pt(3)
        if index < count - 1:
            _rule(slide, 1.2, y + step - 0.12, 11.0, 0.01, _HAIRLINE)
    _add_footer(slide, theme, page_no)


def _layout_for_slide(item: dict[str, Any], title: str, bullets: list[str], index: int, total: int) -> str:
    requested = str(item.get("layout") or "").strip().lower()
    if requested in {"section", "agenda", "cards", "process", "timeline", "comparison", "quote", "summary", "bullets"}:
        return requested

    title_text = str(title or "").lower()
    joined = " ".join(str(bullet or "") for bullet in bullets).lower()
    if re.search(r"总结|结论|建议|下一步|行动|takeaway|summary|recommendation", title_text):
        return "summary"
    if re.search(r"流程|步骤|路径|路线|工作流|计划|roadmap|workflow|process|timeline", title_text):
        return "process"
    if re.search(r"对比|比较|差异|取舍|优劣|vs\.?|versus|compare", title_text + " " + joined):
        return "comparison"
    if len(bullets) <= 2 and (index == 0 or re.search(r"背景|目标|问题|机会|挑战|核心观点|overview|context", title_text)):
        return "quote"
    if 3 <= len(bullets) <= 6:
        return "cards"
    return "bullets"


def _fit_font_size(text: str, *, base: int, small: int) -> int:
    length = len(str(text or ""))
    if length > 90:
        return small
    if length > 58:
        return max(small, base - 3)
    return base


def _style_paragraph(
    paragraph: Any,
    text: str,
    *,
    size: int,
    color: str = _BODY_INK,
    bold: bool = False,
    align: Any | None = None,
) -> None:
    from pptx.util import Pt

    paragraph.text = str(text or "")
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.name = _CN_FONT
    paragraph.font.color.rgb = _hex(color)
    if align is not None:
        paragraph.alignment = align


def _add_header(slide: Any, theme: dict[str, str], title: str) -> None:
    from pptx.util import Inches

    # eyebrow 强调短线（取代写死的英文 kicker chip）；claim 标题用近黑加强层级
    _rule(slide, 0.82, 0.6, 0.62, 0.06, theme["accent"])
    title_box = slide.shapes.add_textbox(Inches(0.82), Inches(0.78), Inches(11.7), Inches(0.92))
    tp = title_box.text_frame.paragraphs[0]
    _style_paragraph(tp, title[:120], size=_fit_font_size(title, base=28, small=22), color=_TITLE_INK, bold=True)


def _add_footer(slide: Any, theme: dict[str, str], page_no: int) -> None:
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches

    footer = slide.shapes.add_textbox(Inches(11.72), Inches(6.92), Inches(0.95), Inches(0.28))
    p = footer.text_frame.paragraphs[0]
    _style_paragraph(p, f"{page_no:02d}", size=9, color=_MUTED_INK, bold=True, align=PP_ALIGN.RIGHT)


def _card(
    slide: Any,
    theme: dict[str, str],
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    body: str = "",
    *,
    number: int | None = None,
) -> None:
    from pptx.util import Inches, Pt

    # 开放式信息块：顶部 accent 短线 +（序号）+ 加粗 lead + 次级 detail，无填充无边框
    _rule(slide, x, y, 0.34, 0.04, theme["accent"])
    box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.16), Inches(w), Inches(h - 0.16))
    tf = box.text_frame
    tf.word_wrap = True
    hp = tf.paragraphs[0]
    lead_size = _fit_font_size(title, base=15, small=12)
    if number is not None:
        nr = hp.add_run()
        nr.text = f"{number:02d}  "
        nr.font.size = Pt(lead_size)
        nr.font.bold = True
        nr.font.name = _CN_FONT
        nr.font.color.rgb = _hex(theme["accent"])
        tr = hp.add_run()
        tr.text = title[:160]
        tr.font.size = Pt(lead_size)
        tr.font.bold = True
        tr.font.name = _CN_FONT
        tr.font.color.rgb = _hex(_TITLE_INK)
    else:
        _style_paragraph(hp, title[:160], size=lead_size, color=_TITLE_INK, bold=True)
    if body:
        bp = tf.add_paragraph()
        _style_paragraph(bp, body[:220], size=_fit_font_size(body, base=11, small=10), color=_DETAIL_INK)
        bp.space_before = Pt(5)


def _split_lead_detail(value: str) -> tuple[str, str]:
    text = str(value or "").strip()
    for delimiter in ("：", ":", " - ", " — ", "——"):
        if delimiter in text:
            left, right = text.split(delimiter, 1)
            if left.strip() and right.strip():
                return left.strip(), right.strip()
    return text, ""


def _add_agenda_slide(prs: Any, theme: dict[str, str], titles: list[str], page_no: int) -> None:
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, "FFFFFF")
    _add_header(slide, theme, "内容导航")

    items = titles[:6]
    y0 = 1.95
    for index, item in enumerate(items, start=1):
        y = y0 + (index - 1) * 0.74
        num = slide.shapes.add_textbox(Inches(0.9), Inches(y - 0.04), Inches(0.8), Inches(0.5))
        _style_paragraph(num.text_frame.paragraphs[0], f"{index:02d}", size=18, color=theme["accent"], bold=True)
        box = slide.shapes.add_textbox(Inches(1.75), Inches(y), Inches(9.8), Inches(0.5))
        _style_paragraph(box.text_frame.paragraphs[0], item[:96], size=17, color=_TITLE_INK, bold=True)
        _rule(slide, 1.75, y + 0.54, 9.6, 0.012, _HAIRLINE)

    _add_footer(slide, theme, page_no)


def _add_cards_slide(prs: Any, theme: dict[str, str], title: str, bullets: list[str], page_no: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, "FFFFFF")
    _add_header(slide, theme, title)
    positions = [
        (0.9, 1.95, 5.5, 1.25),
        (6.95, 1.95, 5.4, 1.25),
        (0.9, 3.35, 5.5, 1.25),
        (6.95, 3.35, 5.4, 1.25),
        (0.9, 4.75, 5.5, 1.25),
        (6.95, 4.75, 5.4, 1.25),
    ]
    for index, bullet in enumerate(bullets[:6]):
        lead, detail = _split_lead_detail(bullet)
        _card(slide, theme, *positions[index], lead, detail, number=index + 1)
    _add_footer(slide, theme, page_no)


def _add_process_slide(prs: Any, theme: dict[str, str], title: str, bullets: list[str], page_no: int) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, "FFFFFF")
    _add_header(slide, theme, title)

    steps = bullets[:5] or ["明确目标", "拆解任务", "执行推进", "验证结果"]
    width = 10.8 / max(1, len(steps))
    y = 3.02
    for index, step in enumerate(steps, start=1):
        x = 0.98 + (index - 1) * width
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + width / 2 - 0.34), Inches(y - 0.32), Inches(0.68), Inches(0.68))
        circle.fill.solid()
        circle.fill.fore_color.rgb = _hex(theme["accent"])
        circle.line.fill.background()
        circle.shadow.inherit = False
        cp = circle.text_frame.paragraphs[0]
        _style_paragraph(cp, str(index), size=13, color="FFFFFF", bold=True, align=PP_ALIGN.CENTER)

        if index < len(steps):
            _rule(slide, x + width / 2 + 0.36, y, max(0.2, width - 0.72), 0.03, _HAIRLINE)

        lead, detail = _split_lead_detail(step)
        box = slide.shapes.add_textbox(Inches(x + 0.08), Inches(y + 0.62), Inches(width - 0.16), Inches(1.25))
        tf = box.text_frame
        tf.word_wrap = True
        _style_paragraph(tf.paragraphs[0], lead[:80], size=_fit_font_size(lead, base=13, small=9), color=_TITLE_INK, bold=True, align=PP_ALIGN.CENTER)
        if detail:
            dp = tf.add_paragraph()
            _style_paragraph(dp, detail[:120], size=9, color=_DETAIL_INK, align=PP_ALIGN.CENTER)
            dp.space_before = Pt(4)
    _add_footer(slide, theme, page_no)


def _add_comparison_slide(prs: Any, theme: dict[str, str], title: str, bullets: list[str], page_no: int) -> None:
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, "FFFFFF")
    _add_header(slide, theme, title)

    # 中线分隔两栏，去掉填充面板，靠对齐与细线建立结构
    _rule(slide, 6.62, 1.95, 0.02, 4.45, _HAIRLINE)
    midpoint = max(1, (len(bullets) + 1) // 2)
    columns = [bullets[:midpoint], bullets[midpoint:] or bullets[:midpoint]]
    headers = ("方案 / 维度 A", "方案 / 维度 B")
    xs = (0.9, 6.95)
    for col, items in enumerate(columns):
        x = xs[col]
        hb = slide.shapes.add_textbox(Inches(x), Inches(1.95), Inches(5.2), Inches(0.4))
        _style_paragraph(hb.text_frame.paragraphs[0], headers[col], size=14, color=theme["accent"], bold=True)
        _rule(slide, x, 2.4, 0.32, 0.04, theme["accent"])
        y = 2.7
        for item in items[:5]:
            lead, detail = _split_lead_detail(item)
            box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(5.3), Inches(0.5))
            box.text_frame.word_wrap = True
            _style_paragraph(box.text_frame.paragraphs[0], lead[:100], size=13, color=_TITLE_INK, bold=True)
            if detail:
                db = slide.shapes.add_textbox(Inches(x), Inches(y + 0.32), Inches(5.3), Inches(0.4))
                db.text_frame.word_wrap = True
                _style_paragraph(db.text_frame.paragraphs[0], detail[:150], size=10, color=_DETAIL_INK)
                y += 0.86
            else:
                y += 0.56
    _add_footer(slide, theme, page_no)


def _add_quote_slide(prs: Any, theme: dict[str, str], title: str, bullets: list[str], page_no: int) -> None:
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, theme["cover_bg"])
    _rule(slide, 0.84, 1.5, 0.1, 4.2, theme["accent"])

    quote = bullets[0] if bullets else title
    box = slide.shapes.add_textbox(Inches(1.25), Inches(1.62), Inches(10.8), Inches(2.6))
    tf = box.text_frame
    tf.word_wrap = True
    _style_paragraph(tf.paragraphs[0], title[:120], size=_fit_font_size(title, base=34, small=26), color="FFFFFF", bold=True)
    qp = tf.add_paragraph()
    _style_paragraph(qp, quote[:220], size=_fit_font_size(quote, base=18, small=13), color=theme["band"])
    qp.space_before = Pt(18)

    # 子点：accent 短线 + 浅色 lead + 更浅 detail（深底上直接用浅色文字，不用 _card）
    for index, bullet in enumerate(bullets[1:4]):
        x = 1.26 + index * 3.7
        lead, detail = _split_lead_detail(bullet)
        _rule(slide, x, 5.18, 0.3, 0.04, theme["accent"])
        b = slide.shapes.add_textbox(Inches(x), Inches(5.32), Inches(3.4), Inches(1.05))
        b.text_frame.word_wrap = True
        _style_paragraph(b.text_frame.paragraphs[0], lead[:80], size=13, color="FFFFFF", bold=True)
        if detail:
            dp = b.text_frame.add_paragraph()
            _style_paragraph(dp, detail[:120], size=10, color=theme["band"])
            dp.space_before = Pt(3)
    _add_footer(slide, theme, page_no)


def _add_summary_slide(prs: Any, theme: dict[str, str], title: str, bullets: list[str], page_no: int) -> None:
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, "FFFFFF")
    _add_header(slide, theme, title)

    # 强调条（直角实心）+ 白色 headline，左对齐更显编排感
    _rule(slide, 0.9, 1.9, 11.5, 1.15, theme["cover_bg"])
    headline = bullets[0] if bullets else "回顾重点，明确下一步行动。"
    hb = slide.shapes.add_textbox(Inches(1.25), Inches(2.12), Inches(10.8), Inches(0.75))
    hb.text_frame.word_wrap = True
    _style_paragraph(hb.text_frame.paragraphs[0], headline[:160], size=_fit_font_size(headline, base=19, small=14), color="FFFFFF", bold=True)

    for index, bullet in enumerate((bullets[1:] or bullets)[:4], start=1):
        lead, detail = _split_lead_detail(bullet)
        x = 1.0 + ((index - 1) % 2) * 5.75
        y = 3.6 + ((index - 1) // 2) * 1.45
        _card(slide, theme, x, y, 5.3, 1.25, lead, detail, number=index)
    _add_footer(slide, theme, page_no)


def _add_rich_slide(
    prs: Any,
    theme: dict[str, str],
    title: str,
    bullets: list[str],
    *,
    layout: str,
    page_no: int,
) -> None:
    if layout == "agenda":
        _add_agenda_slide(prs, theme, bullets or [title], page_no)
    elif layout == "section":
        _add_quote_slide(prs, theme, title, bullets, page_no)
    elif layout == "cards":
        _add_cards_slide(prs, theme, title, bullets, page_no)
    elif layout in {"process", "timeline"}:
        _add_process_slide(prs, theme, title, bullets, page_no)
    elif layout == "comparison":
        _add_comparison_slide(prs, theme, title, bullets, page_no)
    elif layout == "quote":
        _add_quote_slide(prs, theme, title, bullets, page_no)
    elif layout == "summary":
        _add_summary_slide(prs, theme, title, bullets, page_no)
    else:
        _add_content_slide(prs, theme, title, bullets, page_no)


def create_presentation(title: str, slides: Any, *, subtitle: str = "") -> dict[str, Any]:
    """根据大纲生成一份带配色主题的 16:9 .pptx，返回包含 downloadUrl 的结果字典。"""
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ModuleNotFoundError as exc:
        raise AppError(
            "PPT 生成需要 python-pptx，请先安装：pip install python-pptx",
            code=ErrorCode.INVALID_PAYLOAD,
        ) from exc

    title = str(title or "").strip()
    if not title:
        raise AppError("演示文稿需要一个标题（title）。", code=ErrorCode.INVALID_PAYLOAD)
    if not isinstance(slides, list) or not slides:
        raise AppError("演示文稿至少需要一页内容（slides）。", code=ErrorCode.INVALID_PAYLOAD)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    theme = _deck_theme_for(title)

    _add_cover_slide(prs, theme, title, str(subtitle or "").strip())

    normalized_slides: list[dict[str, Any]] = []
    for item in slides[:MAX_SLIDES]:
        if not isinstance(item, dict):
            continue
        slide_title = (str(item.get("title") or "").strip() or f"第 {len(normalized_slides) + 1} 页")[:120]
        bullets = _normalize_bullets(item)[:MAX_BULLETS_PER_SLIDE]
        normalized_slides.append({"raw": item, "title": slide_title, "bullets": bullets})

    if not normalized_slides:
        raise AppError("没有解析到有效的幻灯片内容。", code=ErrorCode.INVALID_PAYLOAD)

    if len(normalized_slides) >= 4:
        _add_agenda_slide(prs, theme, [item["title"] for item in normalized_slides], len(prs.slides) + 1)

    outline: list[dict[str, Any]] = []
    for index, item in enumerate(normalized_slides):
        slide_title = item["title"]
        bullets = item["bullets"]
        layout = _layout_for_slide(item["raw"], slide_title, bullets, index, len(normalized_slides))
        _add_rich_slide(prs, theme, slide_title, bullets, layout=layout, page_no=len(prs.slides) + 1)
        outline.append({"page": index + 1, "title": slide_title, "bullets": bullets, "layout": layout})

    slide_count = len(prs.slides)
    stored = store_generated_file(title, "pptx", lambda path: prs.save(str(path)))

    return {
        "fileId": stored["fileId"],
        "filename": stored["filename"],
        "slideCount": slide_count,
        "downloadUrl": stored["downloadUrl"],
        "title": title,
        "outline": outline,
        "note": (
            "PPT 已生成。请在最终回复里让用户看到制作过程：先一句话说明这份 PPT 的标题和共几页，"
            "再按 outline 字段【逐页展示】每一页的标题和要点（用清晰的分页小标题 + 项目符号列出，不要省略），"
            "最后用 Markdown 链接 [下载 PPT](downloadUrl) 给出下载（链接 6 小时内有效）。"
        ),
    }
