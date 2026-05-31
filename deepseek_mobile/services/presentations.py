"""Generate downloadable .pptx decks from a structured outline.

Exposed to the model as the ``create_pptx`` function-calling tool: the model
passes a title and a list of slides, we render a real PowerPoint file into
``.generated/{id}.pptx`` and hand back a ``downloadUrl`` the model surfaces to
the user as a Markdown link. No frontend changes needed — the chat renderer
already turns Markdown links into clickable anchors.
"""

from __future__ import annotations

import re
import secrets
import shutil
import time
from pathlib import Path
from typing import Any

from deepseek_mobile.core.config import GENERATED_DIR
from deepseek_mobile.core.errors import AppError, ErrorCode
from deepseek_mobile.services.slides_skill import SLIDES_SKILL_NAME

GENERATED_FILE_MAX_AGE_SECONDS = 6 * 3600  # 生成文件保留 6 小时后清理
MAX_SLIDES = 60
MAX_BULLETS_PER_SLIDE = 24
_ID_RE = re.compile(r"[0-9a-f]{32}")
_OUTLINE_HEADING_RE = re.compile(
    r"^\s*(?:#{1,4}\s*)?(?:(?:第\s*)?\d{1,2}\s*[页、.．):：-]|slide\s+\d{1,2}\s*[:：.-])\s*(.+?)\s*$",
    flags=re.IGNORECASE,
)
_BULLET_RE = re.compile(r"^\s*(?:[-*•·]|[（(]?[a-zA-Z][）)]|[（(]?\d{1,2}[）)])\s+(.+?)\s*$")
_REFUSAL_RE = re.compile(r"(?:无法|不能|没有.*能力|只能|可以帮.*大纲|不能直接|无法直接)")


def _safe_filename(title: str) -> str:
    name = re.sub(r"[^\w一-鿿 -]", "", str(title or "")).strip()
    return (name[:60] or "presentation")


def cleanup_generated_files() -> None:
    """删除过期的生成文件。任何 IO 异常都静默吞掉，不影响主流程。"""
    try:
        if not GENERATED_DIR.exists():
            return
        now = time.time()
        for path in GENERATED_DIR.glob("*.pptx"):
            try:
                if now - path.stat().st_mtime > GENERATED_FILE_MAX_AGE_SECONDS:
                    path.unlink(missing_ok=True)
            except OSError:
                continue
    except OSError:
        pass


def resolve_generated_file(file_id: str) -> Path | None:
    """把下载请求的 id 解析成磁盘路径；只接受 32 位十六进制 id，杜绝路径遍历。"""
    if not _ID_RE.fullmatch(str(file_id or "")):
        return None
    path = GENERATED_DIR / f"{file_id}.pptx"
    return path if path.is_file() else None


def save_generated_file_to_downloads(file_id: str, filename: str = "", *, downloads_dir: Path | None = None) -> dict[str, Any]:
    source = resolve_generated_file(file_id)
    if source is None:
        raise AppError("文件不存在或已过期", code=ErrorCode.NOT_FOUND, status=404)

    target_dir = downloads_dir or _downloads_dir()
    target_dir.mkdir(parents=True, exist_ok=True)
    target = _unique_download_path(target_dir, _download_filename(filename))
    shutil.copy2(source, target)
    return {"ok": True, "filename": target.name, "path": str(target), "downloadUrl": f"/api/download?id={file_id}"}


def _downloads_dir() -> Path:
    home = Path.home()
    return home / "Downloads"


def _download_filename(value: str) -> str:
    raw = str(value or "").strip()
    raw = re.sub(r"\.pptx$", "", raw, flags=re.IGNORECASE)
    raw = re.sub(r"^(?:点击)?下载\s*", "", raw, flags=re.IGNORECASE).strip()
    return f"{_safe_filename(raw) or 'presentation'}.pptx"


def _unique_download_path(directory: Path, filename: str) -> Path:
    base = Path(filename).stem or "presentation"
    suffix = Path(filename).suffix or ".pptx"
    candidate = directory / f"{base}{suffix}"
    index = 2
    while candidate.exists():
        candidate = directory / f"{base} ({index}){suffix}"
        index += 1
    return candidate


def _normalize_bullets(item: dict[str, Any]) -> list[str]:
    bullets = item.get("bullets")
    if isinstance(bullets, list):
        return [str(b).strip() for b in bullets if str(b).strip()]
    content = str(item.get("content") or "").strip()
    if content:
        return [line.strip() for line in content.splitlines() if line.strip()]
    return []


def infer_presentation_title(prompt: str, outline: str = "") -> str:
    text = re.sub(r"\s+", " ", str(prompt or "")).strip()
    patterns = (
        r"(?:关于|介绍|讲解|主题为|题为)\s*([A-Za-z0-9\u4e00-\u9fff _-]{1,48}?)\s*(?:的)?\s*(?:PPT|ppt|幻灯片|演示文稿|presentation)",
        r"(?:PPT|ppt|幻灯片|演示文稿|presentation).*?(?:关于|介绍|讲解)\s*([A-Za-z0-9\u4e00-\u9fff _-]{1,48})",
        r"([A-Za-z0-9\u4e00-\u9fff _-]{1,48}?)\s*(?:PPT|ppt|幻灯片|演示文稿|presentation)",
    )
    for pattern in patterns:
        match = re.search(pattern, text, flags=re.IGNORECASE)
        if match:
            title = _clean_outline_title(match.group(1))
            title = re.sub(r"^(?:一个|一份|一套|做一个|做一份|制作一个|制作一份|帮我|请帮我)\s*", "", title)
            if title:
                normalized = _normalize_topic_name(title)
                if re.search(r"\b介绍\b|介绍", text) and "介绍" not in normalized:
                    return f"{normalized} 介绍"
                return normalized

    for line in str(outline or "").splitlines():
        cleaned = _clean_outline_title(line)
        if cleaned and len(cleaned) <= 60 and re.search(r"(?:PPT|幻灯片|演示文稿|介绍)$", cleaned, flags=re.IGNORECASE):
            return _normalize_topic_name(cleaned)
    return "演示文稿"


def create_presentation_from_text(prompt: str, outline: str) -> dict[str, Any]:
    title = infer_presentation_title(prompt, outline)
    slides = slides_from_outline_text(outline, topic=title)
    return create_presentation(title, slides, subtitle=f"由 DeepSeek Mobile {SLIDES_SKILL_NAME} skill 本地生成")


def slides_from_outline_text(outline: str, *, topic: str = "") -> list[dict[str, Any]]:
    text = str(outline or "").replace("\r\n", "\n").replace("\r", "\n")
    text = _outline_relevant_tail(text)
    slides: list[dict[str, Any]] = []
    current: dict[str, Any] | None = None
    for raw_line in text.splitlines():
        line = raw_line.strip()
        if not line or line.startswith("```"):
            continue
        heading = _OUTLINE_HEADING_RE.match(line)
        if heading:
            title = _clean_outline_title(heading.group(1))
            if not title or _REFUSAL_RE.search(title):
                continue
            current = {"title": title, "bullets": []}
            slides.append(current)
            continue
        if current is None:
            continue
        bullet = _BULLET_RE.match(line)
        if bullet:
            text = _clean_outline_bullet(bullet.group(1))
            if text and not _REFUSAL_RE.search(text):
                current["bullets"].append(text)
        elif raw_line.startswith((" ", "\t")):
            text = _clean_outline_bullet(line)
            if text and not _REFUSAL_RE.search(text):
                current["bullets"].append(text)

    slides = _drop_duplicate_cover_slide(slides, topic)
    result: list[dict[str, Any]] = []
    for slide in slides:
        title = _clean_outline_title(slide.get("title", ""))
        if not title:
            continue
        bullets = [str(item).strip() for item in slide.get("bullets", []) if str(item).strip()]
        if not bullets:
            bullets = _default_bullets_for_title(title, topic)
        result.append({"title": title, "bullets": bullets[:MAX_BULLETS_PER_SLIDE]})
        if len(result) >= MAX_SLIDES:
            break
    return result or _default_slides(topic)


def _outline_relevant_tail(text: str) -> str:
    markers = ("PPT 大纲", "ppt 大纲", "幻灯片大纲", "演示文稿大纲", "关于")
    best = -1
    lowered = text.lower()
    for marker in markers:
        index = lowered.rfind(marker.lower())
        if index > best:
            best = index
    return text[best:] if best >= 0 else text


def _clean_outline_title(value: object) -> str:
    text = str(value or "").strip().strip("#*`-—:：.。 ")
    text = re.sub(r"^(?:标题|页面|幻灯片|Slide)\s*[:：-]\s*", "", text, flags=re.IGNORECASE)
    text = re.sub(r"^封面\s*[-:：]\s*", "", text)
    return text[:120].strip()


def _clean_outline_bullet(value: object) -> str:
    text = str(value or "").strip().strip("-*•· ")
    return text[:500].strip()


def _normalize_topic_name(value: str) -> str:
    text = re.sub(r"\s+", " ", str(value or "")).strip(" ：:-")
    if text.lower() == "git":
        return "Git 介绍"
    return text[:80] or "演示文稿"


def _drop_duplicate_cover_slide(slides: list[dict[str, Any]], topic: str) -> list[dict[str, Any]]:
    if len(slides) <= 1:
        return slides
    first_title = str(slides[0].get("title") or "")
    compact_topic = re.sub(r"\s+", "", topic.lower())
    compact_first = re.sub(r"\s+", "", first_title.lower())
    if "封面" in first_title or (compact_topic and compact_first in {compact_topic, f"{compact_topic}介绍"}):
        return slides[1:]
    return slides


def _default_bullets_for_title(title: str, topic: str) -> list[str]:
    topic_text = topic or "本主题"
    return [
        f"说明“{title}”在 {topic_text} 中的作用",
        "列出关键概念、常见场景和易错点",
        "结合简短示例帮助听众快速理解",
    ]


def _default_slides(topic: str) -> list[dict[str, Any]]:
    normalized_topic = topic or "主题"
    if re.search(r"\bgit\b", normalized_topic, flags=re.IGNORECASE):
        return [
            {"title": "什么是 Git", "bullets": ["分布式版本控制系统", "记录代码和文档的每次变更", "支持多人并行协作"]},
            {"title": "为什么需要版本控制", "bullets": ["追踪历史，方便回滚", "多人协作时降低冲突成本", "用分支隔离不同功能和实验"]},
            {"title": "核心概念", "bullets": ["仓库 Repository", "提交 Commit", "分支 Branch", "合并 Merge", "远程仓库 Remote"]},
            {"title": "基本工作流程", "bullets": ["clone 或 init 创建仓库", "add 暂存变更", "commit 保存快照", "push / pull 同步远程"]},
            {"title": "常用命令", "bullets": ["git status 查看状态", "git add / commit 提交变更", "git log 查看历史", "git branch / merge 管理分支"]},
            {"title": "分支与协作", "bullets": ["为功能创建独立分支", "通过 Pull Request / Merge Request 审查", "解决冲突后再合并到主线"]},
            {"title": "最佳实践", "bullets": ["小步提交并写清楚提交信息", "不要提交密钥和大文件", "合并前先同步主分支", "用 .gitignore 排除临时文件"]},
            {"title": "总结", "bullets": ["Git 让变更可追踪、协作可控", "掌握工作流比死记命令更重要", "从日常 add / commit / push 开始练习"]},
        ]
    return [
        {"title": f"{normalized_topic}概览", "bullets": ["背景与目标", "核心问题", "本次演示的主要结论"]},
        {"title": "核心概念", "bullets": ["关键术语", "基本原理", "适用范围"]},
        {"title": "方法与流程", "bullets": ["准备工作", "执行步骤", "检查与反馈"]},
        {"title": "案例或示例", "bullets": ["典型场景", "操作示例", "结果解读"]},
        {"title": "风险与建议", "bullets": ["常见误区", "注意事项", "改进建议"]},
        {"title": "总结", "bullets": ["回顾重点", "下一步行动", "问答"]},
    ]


# 一组现代配色主题；按标题哈希确定性选一套，让同一份内容稳定、不同主题之间有区分。
# 每套：cover_bg 封面深色底，accent 强调色（内容页标题/线/bullet 符号），band 封面副标题浅色。
_DECK_THEMES = (
    {"cover_bg": "1E3A8A", "accent": "2563EB", "band": "BFDBFE"},  # 蓝
    {"cover_bg": "312E81", "accent": "6366F1", "band": "C7D2FE"},  # 靛紫
    {"cover_bg": "0F766E", "accent": "0D9488", "band": "99F6E4"},  # 青绿
    {"cover_bg": "9A3412", "accent": "EA580C", "band": "FED7AA"},  # 暖橙
    {"cover_bg": "0F172A", "accent": "0EA5E9", "band": "BAE6FD"},  # 深空蓝
    {"cover_bg": "581C87", "accent": "A855F7", "band": "E9D5FF"},  # 紫
)
_BODY_INK = "334155"    # 正文深灰
_MUTED_INK = "94A3B8"   # 页码 / 次要
_CN_FONT = "微软雅黑"


def _hex(value: str) -> Any:
    from pptx.dml.color import RGBColor

    return RGBColor.from_string(value)


def _deck_theme_for(title: str) -> dict[str, str]:
    import hashlib

    digest = hashlib.md5(str(title or "").encode("utf-8")).hexdigest()
    return _DECK_THEMES[int(digest, 16) % len(_DECK_THEMES)]


def _set_slide_bg(slide: Any, color: str) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _hex(color)


def _add_cover_slide(prs: Any, theme: dict[str, str], title: str, subtitle: str) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, theme["cover_bg"])

    # 左侧强调竖条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), prs.slide_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = _hex(theme["accent"])
    bar.line.fill.background()
    bar.shadow.inherit = False

    box = slide.shapes.add_textbox(Inches(1.2), Inches(2.5), Inches(10.8), Inches(2.4))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title[:120]
    p.font.size = Pt(40 if len(title) > 18 else 46)
    p.font.bold = True
    p.font.name = _CN_FONT
    p.font.color.rgb = _hex("FFFFFF")
    if subtitle:
        sp = tf.add_paragraph()
        sp.text = subtitle[:160]
        sp.font.size = Pt(18)
        sp.font.name = _CN_FONT
        sp.font.color.rgb = _hex(theme["band"])
        sp.space_before = Pt(16)

    # 标题下方一条短强调线
    underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.25), Inches(4.55), Inches(2.6), Inches(0.08))
    underline.fill.solid()
    underline.fill.fore_color.rgb = _hex(theme["accent"])
    underline.line.fill.background()
    underline.shadow.inherit = False


def _add_content_slide(prs: Any, theme: dict[str, str], title: str, bullets: list[str], page_no: int) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, "FFFFFF")

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.7), Inches(1.0))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title[:120]
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.name = _CN_FONT
    tp.font.color.rgb = _hex(theme["accent"])

    # 标题下强调细线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.85), Inches(1.45), Inches(2.2), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = _hex(theme["accent"])
    line.line.fill.background()
    line.shadow.inherit = False

    # 正文 bullet：数量越多字号越小，避免溢出
    count = max(1, len(bullets))
    size = 22 if count <= 4 else 20 if count <= 6 else 18 if count <= 9 else 16
    body_box = slide.shapes.add_textbox(Inches(0.95), Inches(1.85), Inches(11.4), Inches(5.0))
    btf = body_box.text_frame
    btf.word_wrap = True
    for index, bullet in enumerate(bullets):
        para = btf.paragraphs[0] if index == 0 else btf.add_paragraph()
        marker = para.add_run()
        marker.text = "▸  "
        marker.font.size = Pt(size)
        marker.font.bold = True
        marker.font.name = _CN_FONT
        marker.font.color.rgb = _hex(theme["accent"])
        text_run = para.add_run()
        text_run.text = str(bullet)[:500]
        text_run.font.size = Pt(size)
        text_run.font.name = _CN_FONT
        text_run.font.color.rgb = _hex(_BODY_INK)
        para.space_after = Pt(12)
        try:
            para.line_spacing = 1.15
        except (TypeError, ValueError):
            pass

    # 右下页码
    pn = slide.shapes.add_textbox(Inches(12.2), Inches(6.95), Inches(0.9), Inches(0.4))
    pp = pn.text_frame.paragraphs[0]
    pp.text = str(page_no)
    pp.font.size = Pt(11)
    pp.font.name = _CN_FONT
    pp.font.color.rgb = _hex(_MUTED_INK)


def create_presentation(title: str, slides: Any, *, subtitle: str = "") -> dict[str, Any]:
    """根据大纲生成一份带配色主题的 16:9 .pptx，返回包含 downloadUrl 的结果字典。"""
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ModuleNotFoundError as exc:
        raise AppError(
            "PPT 生成需要 python-pptx，请先安装：pip install python-pptx",
            code=ErrorCode.INVALID_PAYLOAD,
        ) from exc

    title = str(title or "").strip()
    if not title:
        raise AppError("演示文稿需要一个标题（title）。", code=ErrorCode.INVALID_PAYLOAD)
    if not isinstance(slides, list) or not slides:
        raise AppError("演示文稿至少需要一页内容（slides）。", code=ErrorCode.INVALID_PAYLOAD)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    theme = _deck_theme_for(title)

    _add_cover_slide(prs, theme, title, str(subtitle or "").strip())

    content_pages = 0
    outline: list[dict[str, Any]] = []
    for item in slides[:MAX_SLIDES]:
        if not isinstance(item, dict):
            continue
        slide_title = (str(item.get("title") or "").strip() or f"第 {content_pages + 1} 页")[:120]
        bullets = _normalize_bullets(item)[:MAX_BULLETS_PER_SLIDE]
        _add_content_slide(prs, theme, slide_title, bullets, content_pages + 2)
        outline.append({"page": content_pages + 1, "title": slide_title, "bullets": bullets})
        content_pages += 1

    if content_pages == 0:
        raise AppError("没有解析到有效的幻灯片内容。", code=ErrorCode.INVALID_PAYLOAD)

    GENERATED_DIR.mkdir(parents=True, exist_ok=True)
    cleanup_generated_files()
    file_id = secrets.token_hex(16)
    prs.save(str(GENERATED_DIR / f"{file_id}.pptx"))

    return {
        "fileId": file_id,
        "filename": f"{_safe_filename(title)}.pptx",
        "slideCount": content_pages + 1,
        "downloadUrl": f"/api/download?id={file_id}",
        "title": title,
        "outline": outline,
        "note": (
            "PPT 已生成。请在最终回复里让用户看到制作过程：先一句话说明这份 PPT 的标题和共几页，"
            "再按 outline 字段【逐页展示】每一页的标题和要点（用清晰的分页小标题 + 项目符号列出，不要省略），"
            "最后用 Markdown 链接 [下载 PPT](downloadUrl) 给出下载（链接 6 小时内有效）。"
        ),
    }
